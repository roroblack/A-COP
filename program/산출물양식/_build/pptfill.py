# -*- coding: utf-8 -*-
"""발표 템플릿 pptx 를 채우는 공용 도구.

python-pptx 는 기존 파일의 슬라이드를 복제하는 기능이 없다. 템플릿의 내용 슬라이드가
6장뿐이라 우리 발표에는 모자라므로 여기서 직접 복제한다.

★그림 관계(r:embed)는 슬라이드마다 rId 가 따로 매겨진다. 그냥 XML만 복사하면
  새 슬라이드에서 그 rId 가 다른 것을 가리키거나 아예 없어서 그림이 사라진다.
  그래서 관계를 새로 달고 rId 를 바꿔 적는다.
"""
import copy
import re

from PIL import Image
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

EMU = 914400
RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
RT_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"


def dup_slide(prs, index):
    """index 슬라이드를 통째로 복제해 맨 뒤에 붙이고 새 슬라이드를 돌려준다."""
    src = prs.slides[index]
    dst = prs.slides.add_slide(src.slide_layout)
    for shp in list(dst.shapes):
        shp._element.getparent().remove(shp._element)

    # 그림 관계를 새 슬라이드에 달고 옛 rId 를 새 rId 로 바꿀 표를 만든다.
    remap = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype == RT_LAYOUT:
            continue
        if rel.is_external:
            remap[rid] = dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            remap[rid] = dst.part.relate_to(rel.target_part, rel.reltype)

    for shp in src.shapes:
        el = copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in ("embed", "link"):
                key = qn("r:" + attr)
                if key in node.attrib and node.attrib[key] in remap:
                    node.attrib[key] = remap[node.attrib[key]]
        dst.shapes._spTree.append(el)
    return dst


def move_slide(prs, frm, to):
    """슬라이드 순서를 바꾼다. add_slide 는 항상 맨 뒤에 붙기 때문에 필요하다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[frm])
    lst.insert(to, ids[frm])


def drop(shape):
    shape._element.getparent().remove(shape._element)


def _runs(tf):
    return [r for p in tf.paragraphs for r in p.runs]


def set_text(shape, lines, fit=True, size=None):
    """도형의 글을 lines 로 바꾼다. 첫 run 의 서식을 그대로 물려받는다.

    ★fit 이 켜져 있으면 글이 길어진 만큼 글자를 줄인다. 템플릿 상자는 폭이 고정이라
      긴 글을 그대로 넣으면 상자 밖으로 흘러 다른 글을 덮는다.
    """
    if isinstance(lines, str):
        lines = [lines]
    tf = shape.text_frame
    before = max((len(l) for l in tf.text.split("\n")), default=1)

    p0 = tf.paragraphs[0]
    if not p0.runs:
        p0.add_run().text = ""
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0.runs[0].text = lines[0]

    tpl = copy.deepcopy(p0._p)
    for line in lines[1:]:
        newp = copy.deepcopy(tpl)
        tf._txBody.append(newp)
        for t in newp.iter(qn("a:t")):
            t.text = line
            break

    if size is not None:
        for r in _runs(tf):
            r.font.size = Pt(size)
        return
    if not fit:
        return
    after = max(len(l) for l in lines)
    if after <= before:
        return
    ratio = max(0.62, float(before) / after)
    for r in _runs(tf):
        cur = r.font.size
        if cur is None:
            continue
        r.font.size = Emu(int(cur * ratio))


def clear_body(slide, top_in=2.5):
    """머리말과 왼쪽 번호줄만 남기고 본문 도형을 지운다. 그림 한 장을 넣을 자리를 만든다.

    ★왼쪽 밖으로 삐져나온 장식선(x 가 음수)이 있다. 왼쪽 것을 다 남기게 짜면
      이 선이 그림 위를 가로지른다. 그래서 번호줄이 있는 좁은 띠만 남긴다.
    """
    for shp in list(slide.shapes):
        if shp.top is None or shp.left is None:
            continue
        if shp.top / EMU < top_in:
            continue
        if 0.4 <= shp.left / EMU <= 1.2:   # 왼쪽 번호줄 띠
            continue
        drop(shp)


def put_image(slide, path, x=1.94, y=2.95, w=17.1, h=7.2):
    """상자 안에 비율을 지키며 넣는다. 늘리지 않고 남는 쪽을 가운데로 민다."""
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    fw, fh = iw * scale, ih * scale
    slide.shapes.add_picture(path, Emu(int((x + (w - fw) / 2) * EMU)),
                             Emu(int((y + (h - fh) / 2) * EMU)),
                             Emu(int(fw * EMU)), Emu(int(fh * EMU)))


def rail(slide, n):
    """왼쪽 번호줄의 표시점을 n 번(1~6)으로 옮긴다. 목차 6항목 중 어디인지 알린다."""
    ys = {1: 3.07, 2: 4.43, 3: 5.78, 4: 7.12, 5: 8.46, 6: 9.82}
    for shp in slide.shapes:
        if shp.has_text_frame or shp.left is None:
            continue
        if abs(shp.left / EMU - 0.62) < 0.12 and shp.width / EMU < 0.8:
            shp.top = Emu(int(ys[n] * EMU))
            return
    raise ValueError("번호줄 표시점을 못 찾았다")


def find(slide, text, exact=False):
    """글 내용으로 도형을 찾는다. 인덱스로 찾으면 템플릿이 바뀔 때 조용히 엉뚱한 걸 고친다."""
    hits = [s for s in slide.shapes if s.has_text_frame
            and (s.text_frame.text.strip() == text if exact else text in s.text_frame.text)]
    if len(hits) != 1:
        raise ValueError("도형 %r 를 %d 개 찾았다" % (text, len(hits)))
    return hits[0]


def at(slide, x, y, tol=0.15):
    """좌표로 도형을 찾는다. 같은 글이 여러 개일 때 쓴다."""
    hits = [s for s in slide.shapes if s.has_text_frame and s.left is not None
            and abs(s.left / EMU - x) < tol and abs(s.top / EMU - y) < tol]
    if len(hits) != 1:
        raise ValueError("좌표 (%s, %s) 에서 %d 개를 찾았다" % (x, y, len(hits)))
    return hits[0]

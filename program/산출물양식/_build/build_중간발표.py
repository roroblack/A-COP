# -*- coding: utf-8 -*-
"""중간 발표회 발표 자료를 만든다.

템플릿은 다른 팀 예시(AI 모의면접)가 채워진 채로 배포됐다. 디자인은 그대로 두고
글만 우리 것으로 바꾼다. 내용 슬라이드가 6장뿐이라 모자라는 만큼 복제해 늘린다.

★글자 크기를 손으로 지정한다. 템플릿 상자는 폭이 고정인데 예시 글보다 우리 글이
  길다. 그대로 넣으면 상자 밖으로 흘러 옆 글을 덮는다.
  이 템플릿의 한글 한 글자 폭은 글자 크기의 약 0.78배다("최종 프로젝트" 7자가
  86pt 로 6.53인치 상자를 꽉 채운다). 그래서 한 줄에 들어가는 글자 수는
  대략 (상자 폭 인치 x 92 / 글자 크기) 다. 아래 size 값은 이 식으로 정했다.

수치는 지어내지 않는다. 근거는 셋이다.
  - 시장과 격차 수치: team_branch/output/A-COPilot_제출표.xlsx 근거출처
  - DoD 판정과 평가 실행 건수: 2026-08-30 개발 콘솔 실측
  - 일정과 담당: program/plan/A-COP_스프린트_에픽_설계.md
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Emu, Pt

from pptfill import (at, clear_body, drop, dup_slide, find, move_slide,
                     put_image, rail, set_text)

HERE = os.path.dirname(os.path.abspath(__file__))
FORMS = os.path.dirname(HERE)
REPO = os.path.dirname(os.path.dirname(FORMS))
CH = lambda n: os.path.join(REPO, "program", "plan", "diagram", "charts", n)
SH = lambda n: os.path.join(REPO, "program", "plan", "diagram", "screens", n)

SRC = os.path.join(FORMS, "[최종_프로젝트]_중간_발표회_발표(예시)_템플릿.pptx")
OUT = os.path.join(FORMS, "[중간발표] 중간 발표회 발표자료_A-COPilot.pptx")

IN = 914400
CAMP = ["SK네트웍스 ", "Family AI 캠프 32기"]
TOPIC = "대주제명: 다중 에이전트 서빙 기반 고객 피드백 분석 및 맞춤형 응대 자동화 시스템"
PRESENTER = ["6Team A-COPilot", "presenter : 송채영",
             "팀원 : 김지혜 · 서유현 · 정세환", "최상욱 · 최연우"]

prs = Presentation(SRC)
S = prs.slides


def camp(slide, one_line=False):
    """오른쪽 위 기수 표기를 32기로 바꾼다. 템플릿엔 1기와 2기가 섞여 있다."""
    for shp in slide.shapes:
        if shp.has_text_frame and "Family AI 캠프" in shp.text_frame.text:
            set_text(shp, "SK네트웍스 Family AI 캠프 32기" if one_line else CAMP,
                     fit=False)


def topic(slide):
    """대주제명 줄을 넣는다. 예시 주제보다 우리 주제가 길어 상자를 함께 넓힌다."""
    shp = find(slide, "대주제명:")
    shp.width = Emu(int(9.6 * IN))
    set_text(shp, TOPIC, size=17)


def head(slide, part, title, n):
    set_text(at(slide, 1.58, 1.18), part, fit=False)
    set_text(at(slide, 4.94, 1.18), title, fit=False)
    rail(slide, n)
    camp(slide)


# ==================================================== 1. 표지
s = S[0]
# 표지 제목 석 줄. 86pt 로는 "에이전트 플랫폼" 여덟 자가 두 줄로 접힌다.
for (x, y), line in zip([(2.16, 2.78), (6.36, 4.46), (4.34, 6.15)],
                        ["고객 응대를", "끝내는 멀티", "에이전트 플랫폼"]):
    set_text(at(s, x, y), line, size=74)
topic(s)
set_text(at(s, 14.92, 1.29), "SK네트웍스 Family AI 캠프 32기", fit=False)
set_text(at(s, 12.67, 8.68), PRESENTER, size=22)
# 템플릿의 작성 안내문 두 줄은 발표 자료에 남을 글이 아니다.
drop(at(s, 2.53, 9.04))
drop(at(s, 2.53, 9.68))

# ==================================================== 2. 목차
camp(S[1])

# ==================================================== PART 표지 세 장
for i in (2, 7, 9):
    topic(S[i])
    camp(S[i], one_line=True)

# ==================================================== 4. 프로젝트 개요
s = S[3]
head(s, "PART 1", "프로젝트 개요", 1)
set_text(at(s, 3.69, 4.71), "A-COP", fit=False)
set_text(at(s, 3.60, 8.57), "8. 17. - 10. 26.", fit=False)
for (x, y), txt in zip(
        [(10.46, 3.31), (10.44, 5.24), (10.43, 7.17), (10.35, 9.11)],
        ["고객 문의를 상태를 가진 Case 로 관리한다",
         "Team 은 제안만 반환하고 실행은 승인 경로에서 한다",
         "Team 을 추가해도 Core 코드는 바뀌지 않는다",
         "판단과 근거를 함께 기록해 재생하고 감사한다"]):
    set_text(at(s, x, y), txt, size=24)

# ==================================================== 5. 주요 기능
s = S[4]
head(s, "PART 1", "주요 기능", 1)
set_text(at(s, 2.24, 3.31), "Case 생명주기", fit=False)
set_text(at(s, 3.81, 5.93), ["접수부터 종료까지", "12개 상태로 추적한다"], size=28)
set_text(at(s, 2.21, 9.25), "보류와 승인 대기를 상태로 표현", size=30)
set_text(at(s, 11.65, 3.31), "승인 경계", fit=False)
set_text(at(s, 13.08, 5.93), ["고위험 Action 은", "사람이 승인해야 실행된다"], size=26)
set_text(at(s, 11.53, 9.25), "같은 요청 10회에 실제 처리는 1회", size=30)

# ==================================================== 6. 기술스택
s = S[5]
head(s, "PART 1", "기술스택", 2)
# 셋째 칸만 상자가 낮다. 셋 다 같은 분량이 들어가므로 높이를 맞춘다.
at(s, 14.12, 5.99).height = Emu(int(3.40 * IN))
set_text(at(s, 2.18, 3.31), "서버", fit=False)
set_text(at(s, 8.18, 3.31), "배포", fit=False)
set_text(at(s, 14.08, 3.31), "DB", fit=False)
set_text(at(s, 2.26, 5.99),
         ["Python, FastAPI, LangGraph", "",
          "최상위 흐름은 LangGraph 로 짜고 Agentic Controller 가 "
          "어느 Team 에 일을 줄지 정한다. Team 은 Registry 등록형이다."], size=20)
set_text(at(s, 8.18, 5.99),
         ["자체 호스팅", "",
          "고객사 환경에 올려 데이터가 밖으로 나가지 않게 한다. "
          "Docker 와 AWS 는 Phase 2 다. 지금 로컬은 PostgreSQL 을 직접 띄운다."], size=20)
set_text(at(s, 14.12, 5.99),
         ["PostgreSQL, pgvector", "",
          "업무 상태와 지식 검색을 한 저장소에 둔다. "
          "MySQL 이 아니라 PostgreSQL 인 이유는 pgvector 확장 때문이다."], size=20)

# ==================================================== 7. 팀원분담
s = S[6]
head(s, "PART 1", "팀원분담", 2)
set_text(at(s, 2.25, 3.15), "Core", fit=False)
set_text(at(s, 11.08, 3.15), "Agent Team Module 과 검증", size=30)
for (x, y), name in zip([(2.28, 5.35), (6.56, 5.35), (10.94, 5.35), (15.26, 5.35)],
                        ["최연우", "정세환", "송채영", "최상욱"]):
    set_text(at(s, x, y), name, fit=False)
for (x, y), lines in zip(
        [(2.43, 7.01), (6.94, 7.01), (11.00, 7.01), (15.49, 7.01)],
        [["· 코어 1", "· Case 생명주기", "· Shared State", "· Controller", "· Registry"],
         ["· 코어 2", "· REST 와 MCP", "· A2A Adapter", "· Action 실행", "· 승인과 감사"],
         ["· 모델 총괄", "· VOC 와 응대", "· 주문결제 김지혜", "· 배송반품 서유현",
          "· 카탈로그 공동"],
         ["· 평가 harness", "· golden 관리", "· holdout 관리", "· 통계 검정",
          "· 운영 콘솔"]]):
    set_text(at(s, x, y), lines, size=19)

# ==================================================== 9. 개발 일정
s = S[8]
head(s, "PART 2", "개발 일정", 3)
set_text(at(s, 7.97, 9.50), "스프린트", fit=False)
for (x, y), txt in zip([(7.81, 5.12), (8.69, 6.60), (12.03, 8.03), (14.62, 9.50)],
                       [["S1 기반 계약 고정", "08.28 - 09.15"],
                        ["S2 도메인 팀 구현", "09.16 - 09.30"],
                        ["S3 평가와 통합", "10.01 - 10.14"],
                        ["S4 배포와 시연", "10.15 - 10.26"]]):
    set_text(at(s, x, y), txt, size=18)

# ==================================================== 11, 12. Q&A 와 마무리
camp(S[10])
camp(S[11], one_line=True)
topic(S[11])
set_text(find(S[11], "1Team"), PRESENTER, size=22)

# ==================================================== 새로 만드는 슬라이드
# 기술스택 슬라이드(3칸)를 원형으로 쓴다. 제목 자리와 본문 자리 좌표가 같다.
COL3 = [(2.18, 3.31, 2.26, 5.99), (8.18, 3.31, 8.18, 5.99), (14.08, 3.31, 14.12, 5.99)]


def three(part, title, n, cols):
    s = dup_slide(prs, 5)
    head(s, part, title, n)
    for (tx, ty, bx, by), (t, body) in zip(COL3, cols):
        set_text(at(s, tx, ty), t, size=24)
        set_text(at(s, bx, by), body, size=20)
    return s


def picture(part, title, n, path, note):
    s = dup_slide(prs, 5)
    head(s, part, title, n)
    clear_body(s)
    put_image(s, path, y=2.95, h=6.4)
    box = s.shapes.add_textbox(Emu(int(1.94 * IN)), Emu(int(9.55 * IN)),
                               Emu(int(17.1 * IN)), Emu(int(0.7 * IN)))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = note
    r.font.size = Pt(16)
    return s


N = len(prs.slides._sldIdLst)

three("PART 1", "왜 필요한가", 1, [
    ("도입은 됐는데 작동하지 않는다",
     "콜센터의 88퍼센트가 AI 를 쓰지만 완전히 통합한 곳은 25퍼센트다. "
     "Gartner 는 2027년까지 agentic AI 프로젝트의 40퍼센트 이상이 폐기될 것으로 봤고, "
     "원인을 모델 실패가 아니라 운영화 실패로 지목했다."),
    ("한 줄 흐름으로는 안 되는 일들",
     "수집하고 분류하고 검색해 답변하는 한 줄 파이프라인으로는 실제 고객센터를 담지 못한다. "
     "며칠 멈춰 있는 보류 건, 사람 결재가 필요한 환불 건, 여러 부서를 거치는 건이 그 밖에 있다."),
    ("1군 제품의 공백",
     "Intercom Fin 은 단일 에이전트이고 자체 호스팅이 되지 않는다는 공개 근거가 있다. "
     "데이터 반출이 금지된 규제 산업에는 그대로 들어가기 어렵다."),
])
picture("PART 1", "도입과 실제의 격차", 1, CH("02_adoption_gap.png"),
        "출처: 제출표 근거 D1, D3. 부족한 것은 모델 성능이 아니라 운영에 얹을 수 있는 구조다.")
picture("PART 1", "시스템 구조", 2, CH("uml_component_v2.png"),
        "파란 영역이 코어 1 Case Runtime, 붉은 영역이 코어 2 Access and Action, "
        "초록 영역이 Registry 에 등록되는 도메인 팩이다.")
three("PART 2", "지금까지 된 것", 3, [
    ("DoD 판정",
     "통과 26, 부분통과 3, 판정 없음 1. 2026-08-30 개발 콘솔로 evidence 폴더를 읽어 판정했다. "
     "막힌 3건은 AB holdout, 마일스톤 게이트, 파인튜닝 방어지표다."),
    ("평가 실행",
     "리포트 40건이 쌓였다. run 단위로만 비교한다. arm 과 dataset 과 실행 방식이 다르면 "
     "평균을 견줄 수 없다. 이 경고를 화면 머리에 고정해 두었다."),
    ("Composer 쓰기 채널",
     "카탈로그 조회와 변경 요청 계약을 참고 구현체에서 먼저 만들어 검증했다. "
     "릴리스 대상으로 이식하는 일이 남아 있다."),
])
picture("PART 2", "무엇이 막혀 있나", 3, SH("shot_scr02_card2_dod.png"),
        "운영 콘솔의 실제 화면이다. 막힌 항목만 표에 남기고 통과한 26건은 접어 둔다.")
three("PART 2", "이슈와 트러블 슈팅", 4, [
    ("평가 점수가 0퍼센트로 나왔다",
     "A 군 정확도가 0.0퍼센트였다. 채점기가 승인 대기를 실패로 세고, 성공 판정을 뒤에서 "
     "덮어쓰고 있었다. 모델이 아니라 채점 코드의 결함이었다. 지시서로 담당에게 넘겼다."),
    ("설정을 고객 서버에서 고치고 있었다",
     "쓰기 가능한 구성 관리가 고객 대면 경로에 있으면 취약점 하나가 서비스 전체에 닿는다. "
     "중앙 설정 저장소로 옮기기로 2026-08-29 에 결정했다."),
    ("적용 버튼이 항상 거부됐다",
     "화면이 사유를 안 보내서 422 였다. 대상이 사유를 필수로 요구한다. "
     "감사 기록의 근거이기 때문이다. 가짜 대상을 세워 끝까지 다시 확인했다."),
])
three("PART 3", "예상 결과 목표", 5, [
    ("숫자로 증명한다",
     "golden 과 holdout 으로 A 군, B 군, 제안군을 비교하고 bootstrap 과 McNemar 로 "
     "검정한다. 좋아졌다를 느낌이 아니라 수치로 말한다."),
    ("멱등성과 승인 경계",
     "같은 요청 10회에 실제 부수효과는 1회다. 고위험 Action 은 승인 없이 실행되지 않는다. "
     "이 둘을 테스트로 고정한다."),
    ("구성으로 늘어나는 확장",
     "Team 을 하나 추가할 때 Core 코드가 바뀌지 않는 것을 아키텍처 테스트로 막는다. "
     "최종 시연은 검증 쇼핑몰 연동이다."),
])

# 새로 만든 일곱 장을 제자리로 옮긴다. add_slide 는 맨 뒤에만 붙기 때문이다.
# 앞에 한 장을 끼울 때마다 뒤가 한 칸씩 밀리므로 옮길 대상 번호도 하나씩 올라간다.
for k, dest in enumerate([5, 6, 8, 12, 13, 14, 16]):
    move_slide(prs, N + k, dest)

prs.save(OUT)
print("저장:", OUT)
print("슬라이드 %d장" % len(prs.slides))
for i, sl in enumerate(prs.slides, 1):
    t = [sh.text_frame.text.replace("\n", " ") for sh in sl.shapes
         if sh.has_text_frame and sh.left is not None
         and abs(sh.left / IN - 4.94) < 0.2 and abs(sh.top / IN - 1.18) < 0.2]
    print("  %2d  %s" % (i, t[0] if t else ""))
